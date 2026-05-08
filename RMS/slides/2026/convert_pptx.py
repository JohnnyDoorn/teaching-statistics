#!/usr/bin/env python3
"""Convert .pptx files to markdown with formatting annotations and extracted images."""
import sys, os, zipfile, hashlib
from pptx import Presentation
from pptx.util import Pt
from pptx.enum.shapes import MSO_SHAPE_TYPE

def rgb_hex(color):
    try:
        return f"#{color}"
    except:
        return None

def extract_images(filepath, img_dir):
    """Extract all images from the pptx and return a mapping of rId -> relative path."""
    os.makedirs(img_dir, exist_ok=True)
    image_map = {}  # (slide_idx, rId) -> relative path
    
    with zipfile.ZipFile(filepath) as z:
        # Extract all media files
        media_files = {}
        for name in z.namelist():
            if name.startswith("ppt/media/"):
                basename = os.path.basename(name)
                outpath = os.path.join(img_dir, basename)
                with z.open(name) as src, open(outpath, "wb") as dst:
                    dst.write(src.read())
                media_files[name] = os.path.join(os.path.basename(img_dir), basename)
    
    return media_files

def get_image_path(shape, slide, media_files, slide_idx):
    """Try to resolve an image shape to its extracted file path."""
    try:
        image = shape.image
        # The image blob's content hash or the partname
        partname = image.content_type  # not useful directly
        # Try via the relationship
        for rel in slide.part.rels.values():
            if hasattr(rel, 'target_partname') and rel.target_partname:
                target = str(rel.target_partname)
                if target.startswith("/"):
                    target = target[1:]
                if target in media_files and rel.reltype and "image" in rel.reltype:
                    # Match by checking if this rel points to our shape's image
                    pass
        # Simpler: use shape.image directly
        blob = image.blob
        ext = image.content_type.split("/")[-1]
        if ext == "jpeg": ext = "jpg"
        if ext == "x-emf": ext = "emf"
        if ext == "x-wmf": ext = "wmf"
        # Find matching file by content
        img_dir = os.path.dirname(list(media_files.values())[0]) if media_files else "img"
        for media_path, rel_path in media_files.items():
            full_path = rel_path  # relative path
            # Check parent dir for the file
            if os.path.exists(rel_path):
                with open(rel_path, "rb") as f:
                    if f.read() == blob:
                        return rel_path
        # Fallback: save with shape name
        safe_name = "".join(c if c.isalnum() or c in "-_" else "_" for c in shape.name)
        fname = f"{safe_name}.{ext}"
        fpath = os.path.join(img_dir, fname)
        with open(fpath, "wb") as f:
            f.write(blob)
        return fpath
    except Exception:
        return None

def extract_shape(shape, indent="", slide=None, media_files=None, slide_idx=0):
    lines = []
    st = shape.shape_type

    if shape.has_text_frame:
        for para in shape.text_frame.paragraphs:
            text = para.text.strip()
            if not text:
                continue
            max_size = 0
            fmt_parts = []
            for run in para.runs:
                f = run.font
                p = []
                if f.bold: p.append("bold")
                if f.italic: p.append("italic")
                if f.size:
                    p.append(f"{f.size.pt:.0f}pt")
                    max_size = max(max_size, f.size.pt)
                try:
                    if f.color and f.color.rgb:
                        p.append(rgb_hex(f.color.rgb))
                except:
                    pass
                if p:
                    fmt_parts.append("[" + ", ".join(p) + "]")

            if max_size >= 28:
                line = f"{indent}## {text}"
            elif max_size >= 20:
                line = f"{indent}### {text}"
            else:
                prefix = "  " * para.level
                line = f"{indent}{prefix}- {text}"

            if fmt_parts:
                line += f"  <!-- {' '.join(fmt_parts)} -->"
            lines.append(line)

    if st == MSO_SHAPE_TYPE.PICTURE:
        img_path = get_image_path(shape, slide, media_files, slide_idx) if slide else None
        if img_path:
            lines.append(f"{indent}![{shape.name}]({img_path})")
        else:
            lines.append(f"{indent}![Image: {shape.name}]()")
    elif st == MSO_SHAPE_TYPE.TABLE:
        table = shape.table
        rows = list(table.rows)
        header = "| " + " | ".join(c.text for c in rows[0].cells) + " |"
        sep = "| " + " | ".join("---" for _ in rows[0].cells) + " |"
        lines.append(indent + header)
        lines.append(indent + sep)
        for row in rows[1:]:
            lines.append(indent + "| " + " | ".join(c.text for c in row.cells) + " |")
    elif st == MSO_SHAPE_TYPE.GROUP:
        lines.append(f"{indent}<!-- Group: {shape.name} -->")
        for s in shape.shapes:
            lines.extend(extract_shape(s, indent + "  ", slide, media_files, slide_idx))
    elif not shape.has_text_frame:
        stype = str(st).split(".")[-1] if st else "unknown"
        lines.append(f"{indent}<!-- Shape: {shape.name} ({stype}) -->")

    return lines

def convert(filepath):
    prs = Presentation(filepath)
    name = os.path.splitext(os.path.basename(filepath))[0]
    img_dir = name + "_images"
    media_files = extract_images(filepath, img_dir)
    out = [f"# {name}\n"]

    for i, slide in enumerate(prs.slides, 1):
        out.append(f"\n---\n\n## Slide {i}\n")
        try:
            out.append(f"<!-- Layout: {slide.slide_layout.name} -->\n")
        except:
            pass

        for shape in slide.shapes:
            info = extract_shape(shape, slide=slide, media_files=media_files, slide_idx=i)
            if info:
                out.extend(info)
                out.append("")

        if slide.has_notes_slide and slide.notes_slide.notes_text_frame:
            notes = slide.notes_slide.notes_text_frame.text.strip()
            if notes:
                out.append(f"\n> **Notes:** {notes}\n")

    return "\n".join(out)

if __name__ == "__main__":
    for f in sys.argv[1:]:
        md = convert(f)
        outname = os.path.splitext(f)[0] + ".md"
        with open(outname, "w") as fh:
            fh.write(md)
        nslides = md.count("## Slide ")
        print(f"{f} -> {outname} ({nslides} slides, {len(md)} chars)")
